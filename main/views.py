from django.shortcuts import render, get_object_or_404, redirect
from django.contrib.auth.decorators import login_required, user_passes_test
from django.contrib.admin.views.decorators import staff_member_required
from main.utils import can_access, get_all_subordinates
from django.http import HttpResponseForbidden
from django.contrib.auth.models import User
from .models import Level4Text, Profile  # Import the model
from django.http import HttpResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from django.http import HttpResponseRedirect,HttpResponseBadRequest
from django.urls import reverse
from django.conf import settings
from pptx.util import Inches
from io import BytesIO
from .models import Level4Text, Profile
from .utils import get_all_subordinates
from django.views.decorators.http import require_POST
from django.http import JsonResponse

import os
import datetime
# Check if the user is Level 4
def is_level_4(user):
    return user.profile.level == 4


@login_required
def main_page(request):
    user_profile = request.user.profile  # Get the user's profile
    user_level = user_profile.level      # Get the user's level
    user_div = user_profile.div          # Get the user's division
    user_sub_div = user_profile.sub_div  # Get the user's subdivision

    if user_level == 4:
        # Level 4 users see their own highlights and lowlights
        return render(request, 'main.html', {
            'user_level': user_level,
            'user_div': user_div,
            'user_sub_div': user_sub_div,
        })

    elif user_level == 3:
        # Level 3 users see subdivisions directly under them
        subordinates = Profile.objects.filter(manager=user_profile)
        sub_divisions = subordinates.values_list('sub_div', flat=True).distinct()
        return render(request, 'main.html', {
            'user_level': user_level,
            'user_div': user_div,
            'sub_divisions': sub_divisions,
        })

    elif user_level == 2:
        # Level 2 users see aggregated data for all subdivisions under their division
        subordinates = Profile.objects.filter(manager=user_profile)
        sub_divisions = subordinates.values_list('sub_div', flat=True).distinct()
        return render(request, 'main.html', {
            'user_level': user_level,
            'user_div': user_div,
            'sub_divisions': sub_divisions,
        })

    elif user_level == 1:
        # Level 1 users see the admin dashboard
        return render(request, 'main.html', {
            'user_level': user_level,
        })

    else:
        # Redirect to login if the user level is invalid
        return HttpResponseRedirect(reverse('login'))
    
@login_required
def redirect_after_login(request):
    level = request.user.profile.level
    if level == 4:
        return redirect('HL_LL')
    elif level == 3:
        return redirect('level_3_page')
    elif level == 2:
        return redirect('level_2_page')
    elif level == 1:
        return redirect('level_1_page')
    else:
        return redirect('main')

@login_required
def level_4_page(request):
    user_profile = request.user.profile  # Get the user's profile
    user_div = user_profile.div          # Get the user's division
    user_sub_div = user_profile.sub_div  # Get the user's subdivision

    if request.method == 'POST':
        # Handle the submitted highlights and lowlights
        highlight_text = request.POST.get('highlight_text', '').strip()
        lowlight_text = request.POST.get('lowlight_text', '').strip()

        if highlight_text:  # Save highlights if not empty
            Level4Text.objects.create(
                user=request.user,
                sub_div=user_sub_div,
                category='highlight',
                text=highlight_text
            )

        if lowlight_text:  # Save lowlights if not empty
            Level4Text.objects.create(
                user=request.user,
                sub_div=user_sub_div,
                category='lowlight',
                text=lowlight_text
            )

        # Redirect to the same page to prevent form resubmission
        return HttpResponseRedirect(reverse('HL_LL'))

    # Retrieve only the texts written by the current user for their division and subdivision
    all_highlights = Level4Text.objects.filter(
        user=request.user,
        sub_div=user_sub_div,
        category='highlight'
    ).order_by('-created_at')

    all_lowlights = Level4Text.objects.filter(
        user=request.user,
        sub_div=user_sub_div,
        category='lowlight'
    ).order_by('-created_at')

    current_week = datetime.datetime.now().isocalendar()[1]  # Get the current week number

    current_week_highlights = all_highlights.filter(created_at__week=current_week)
    older_highlights = all_highlights.exclude(created_at__week=current_week)

    current_week_lowlights = all_lowlights.filter(created_at__week=current_week)
    older_lowlights = all_lowlights.exclude(created_at__week=current_week)

    return render(request, 'level_4_page.html', {
        'current_week': current_week,
        'current_week_highlights': current_week_highlights,
        'older_highlights': older_highlights,
        'current_week_lowlights': current_week_lowlights,
        'older_lowlights': older_lowlights,
        'user_div': user_div,
        'user_sub_div': user_sub_div,
    })

def level_1_page(request):
    if request.user.profile.level != 1:
        return HttpResponseForbidden("You do not have permission to access this page.")

    # Retrieve all highlights and lowlights grouped by division and subdivision
    divisions = Profile.objects.values_list('div', flat=True).distinct()
    grouped_data_current_week = {}
    grouped_data_old = {}

    current_week = datetime.datetime.now().isocalendar()[1]  # Get the current week number

    for div in divisions:
        subdivisions = Profile.objects.filter(div=div).values_list('sub_div', flat=True).distinct()
        grouped_data_current_week[div] = {}
        grouped_data_old[div] = {}
        for sub_div in subdivisions:
            highlights_current_week = Level4Text.objects.filter(
                sub_div=sub_div, category='highlight', created_at__week=current_week
            ).order_by('-created_at')
            lowlights_current_week = Level4Text.objects.filter(
                sub_div=sub_div, category='lowlight', created_at__week=current_week
            ).order_by('-created_at')

            highlights_old = Level4Text.objects.filter(
                sub_div=sub_div, category='highlight'
            ).exclude(created_at__week=current_week).order_by('-created_at')
            lowlights_old = Level4Text.objects.filter(
                sub_div=sub_div, category='lowlight'
            ).exclude(created_at__week=current_week).order_by('-created_at')

            grouped_data_current_week[div][sub_div] = {
                'highlights': highlights_current_week,
                'lowlights': lowlights_current_week,
            }
            grouped_data_old[div][sub_div] = {
                'highlights': highlights_old,
                'lowlights': lowlights_old,
            }

    return render(request, 'level_1_page.html', {
        'grouped_data_current_week': grouped_data_current_week,
        'grouped_data_old': grouped_data_old,
    })

@login_required
def level_3_page(request):
    if request.user.profile.level != 3:
        return HttpResponseForbidden("You do not have permission to access this page.")

    # Fetch all subordinates recursively
    all_subordinates = get_all_subordinates(request.user.profile)
    # Include the user himself
    all_subordinates = list(all_subordinates) + [request.user.profile]
    subordinate_user_ids = [subordinate.user.id for subordinate in all_subordinates]
    subdivisions = [subordinate.sub_div for subordinate in all_subordinates]
    subdivisions = list(set(subdivisions))  # Remove duplicates

    all_highlights = Level4Text.objects.filter(user__in=subordinate_user_ids, category='highlight').order_by('-created_at')
    all_lowlights = Level4Text.objects.filter(user__in=subordinate_user_ids, category='lowlight').order_by('-created_at')

    current_week = datetime.datetime.now().isocalendar()[1]  # Get the current week number

    current_week_highlights = all_highlights.filter(created_at__week=current_week)
    older_highlights = all_highlights.exclude(created_at__week=current_week)

    current_week_lowlights = all_lowlights.filter(created_at__week=current_week)
    older_lowlights = all_lowlights.exclude(created_at__week=current_week)


    return render(request, 'level_3_page.html', {
        'current_week_highlights': current_week_highlights,
        'older_highlights': older_highlights,
        'current_week_lowlights': current_week_lowlights,
        'older_lowlights': older_lowlights,
        'subdivisions': subdivisions,
    })


@login_required
def level_2_page(request):
    if request.user.profile.level != 2:
        return HttpResponseForbidden("You do not have permission to access this page.")


    # Fetch all subordinates recursively
    all_subordinates = get_all_subordinates(request.user.profile)
    # Include the user himself
    all_subordinates = list(all_subordinates) + [request.user.profile]
    subordinate_user_ids = [subordinate.user.id for subordinate in all_subordinates]
    subdivisions = [subordinate.sub_div for subordinate in all_subordinates]
    subdivisions = list(set(subdivisions))  # Remove duplicates

    # Retrieve highlights and lowlights for all subordinates
    all_highlights = Level4Text.objects.filter(user__in=subordinate_user_ids, category='highlight').order_by('-sub_div', '-created_at')
    all_lowlights = Level4Text.objects.filter(user__in=subordinate_user_ids, category='lowlight').order_by('-sub_div', '-created_at')

    print("Highlights sub_divs:", list(all_highlights.values_list('sub_div', flat=True).distinct()))
    
    current_week = datetime.datetime.now().isocalendar()[1]  # Get the current week number

    current_week_highlights = all_highlights.filter(created_at__week=current_week)
    older_highlights = all_highlights.exclude(created_at__week=current_week)

    current_week_lowlights = all_lowlights.filter(created_at__week=current_week)
    older_lowlights = all_lowlights.exclude(created_at__week=current_week)



    return render(request, 'level_2_page.html', {
        'current_week_highlights': current_week_highlights,
        'older_highlights': older_highlights,
        'current_week_lowlights': current_week_lowlights,
        'older_lowlights': older_lowlights,
        'subdivisions': subdivisions,
    })



@require_POST
@login_required
def toggle_export_selection(request, text_id):
    text = get_object_or_404(Level4Text, id=text_id)
    # Only allow the user's manager or the user to change selection
    if request.user == text.user or request.user.profile.level < text.user.profile.level:
        selected = request.POST.get('selected') == 'true'
        text.export_selected = selected
        text.save()
        return JsonResponse({'success': True, 'selected': text.export_selected})
    return JsonResponse({'success': False}, status=403)

@login_required
def delete_text(request, text_id):
    text = get_object_or_404(Level4Text, id=text_id)
    text.delete()
    if(request.user.profile.level == 4):
        return redirect('HL_LL')
    else:
        return HttpResponseRedirect(request.META.get('HTTP_REFERER', '/'))

@login_required
def edit_text(request, text_id):
    text = get_object_or_404(Level4Text, id=text_id)

    if request.method == 'POST':
        new_text = request.POST.get('typed_text', '')
        text.text = new_text
        text.save()
        return HttpResponseRedirect(request.META.get('HTTP_REFERER', '/'))
    return render(request, 'edit_text.html', {'text': text})

@login_required
def add_text(request):

    if request.method == 'POST':
        sub_div = request.POST.get('sub_div')
        highlight_text = request.POST.get('highlight_text', '').strip()
        lowlight_text = request.POST.get('lowlight_text', '').strip()

        if not sub_div:
            return HttpResponseBadRequest("Subdivision is required.")

        # Add highlight if provided
        if highlight_text:
            Level4Text.objects.create(
                user=request.user,
                sub_div=sub_div,
                category='highlight',
                text=highlight_text
            )

        # Add lowlight if provided
        if lowlight_text:
            Level4Text.objects.create(
                user=request.user,
                sub_div=sub_div,
                category='lowlight',
                text=lowlight_text
            )
        return HttpResponseRedirect(request.META.get('HTTP_REFERER', '/'))

@login_required
def generate_ppt(request):
    # Load your uploaded template
    template_path = os.path.join(settings.BASE_DIR, 'required.pptx')
    presentation = Presentation(template_path)
    

    today = datetime.datetime.today().strftime('%d %b %Y')
    user = request.user
    full_name = str(user.get_full_name() or user.username)
    department = user.profile.sub_div or "Unknown Dept"

    #SLIDE1 - Title
    slide1 = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide1.shapes.title.text = f"{full_name} – Highlights & Lowlights"
    if len(slide1.placeholders) > 1:
        slide1.placeholders[1].text = f"Department: {department}\nDate: {today}"

    selected_highlights_ids = request.POST.getlist('selected_highlights')
    selected_lowlights_ids = request.POST.getlist('selected_lowlights')

    print('selected_highlights', selected_highlights_ids)

    #Get HL/LL data
    if user.profile.level == 4:
        all_highlights = Level4Text.objects.filter(
            id__in=selected_highlights_ids, user=user, category='highlight', export_selected=True
        ).order_by('-created_at')
        all_lowlights = Level4Text.objects.filter(
            id__in=selected_lowlights_ids, user=user, category='lowlight', export_selected=True
        ).order_by('-created_at')
    elif user.profile.level == 1:
        all_highlights = Level4Text.objects.filter(
            id__in=selected_highlights_ids, category='highlight', export_selected=True
        ).order_by('-created_at')
        all_lowlights = Level4Text.objects.filter(
            id__in=selected_lowlights_ids, category='lowlight', export_selected=True
        ).order_by('-created_at')
    else:
        all_subordinates = get_all_subordinates(user.profile) + [user.profile]
        user_ids = [sub.user.id for sub in all_subordinates]
        all_highlights = Level4Text.objects.filter(
            id__in=selected_highlights_ids, user__in=user_ids, category='highlight', export_selected=True
        ).order_by('-created_at')
        all_lowlights = Level4Text.objects.filter(
            id__in=selected_lowlights_ids, user__in=user_ids, category='lowlight', export_selected=True
        ).order_by('-created_at')

    highlights = []
    for t in all_highlights:
        clean_text = t.text.replace('\r', '').replace('\n', ' ')
        highlights.append(f"[{t.sub_div}] {clean_text}")

    lowlights = []
    for t in all_lowlights:
        clean_text = t.text.replace('\r', '').replace('\n', ' ')
        lowlights.append(f"[{t.sub_div}] {clean_text}")

    if not highlights and not lowlights:
        highlights = ["No highlights available."]
        lowlights = ["No lowlights available."]

    # SLIDE 2 — HL & LL Table
    slide2 = presentation.slides.add_slide(presentation.slide_layouts[15])
    #slide2.shapes.title.text = f"{full_name} – HL & LL – {department}"

    row_count = max(1, max(len(highlights), len(lowlights))) + 1
    col_count = 2
    left = Inches(0.8)
    top = Inches(1)
    width = Inches(11.76)
    height = Inches(0.6 * row_count)

    table = slide2.shapes.add_table(row_count, col_count, left, top, width, height).table
    table.cell(0, 0).text = "Highlights"
    table.cell(0, 1).text = "Lowlights"

    for i in range(2):
        cell = table.cell(0, i)
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.name = 'Arial'
        para.font.size = Pt(14)
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(10, 130, 118)

    for i in range(1, row_count):
        table.cell(i, 0).text = highlights[i - 1] if i - 1 < len(highlights) else ""
        table.cell(i, 1).text = lowlights[i - 1] if i - 1 < len(lowlights) else ""
        for j in range(2):
            para = table.cell(i, j).text_frame.paragraphs[0]
            para.font.name = 'Arial'
            para.font.size = Pt(12)
            para.alignment = PP_ALIGN.LEFT
    
    comment_left = Inches(0.8)
    comment_top = Inches(6)
    comment_width = Inches(11.76)
    comment_height = Inches(0.6)
    comment_box = slide2.shapes.add_textbox(comment_left, comment_top, comment_width, comment_height)
    comment_frame = comment_box.text_frame
    comment_frame.text = "Comments:"
    comment_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    comment_frame.paragraphs[0].font.size = Pt(14)
    comment_frame.paragraphs[0].font.bold = True
    comment_frame.paragraphs[0].font.name = 'Arial'
    comment_box.line.color.rgb = RGBColor(10, 130, 118)
    comment_box.line.width = Pt(3)

    # SLIDE 3 — Static Slide
    if len(presentation.slide_layouts) > 2:
       presentation.slides.add_slide(presentation.slide_layouts[17])

    # Return as downloadable PPTX
    ppt_stream = BytesIO()
    presentation.save(ppt_stream)
    ppt_stream.seek(0)

    response = HttpResponse(ppt_stream, content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response['Content-Disposition'] = 'attachment; filename="Highlights_Lowlights.pptx"'
    return response

@staff_member_required
def admin_delete_texts(request):
    if request.method == 'POST':
        deleted_count, _ = Level4Text.objects.all().delete()
        return render(request, 'admin_delete_texts.html', {'deleted_count': deleted_count})
    return render(request, 'admin_delete_texts.html')